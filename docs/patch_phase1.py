# Phase 1 surgical fixes: v5 -> v6. Self-verifying: prints before/after for every change.
import sys
from pptx import Presentation
from pptx.util import Emu, Pt, Inches

sys.stdout.reconfigure(encoding="utf-8")

SRC = r"E:\FARAWAY\docs\KAVACH_FARAWAY2026_FINAL_v5.pptx"
DST = r"E:\FARAWAY\docs\KAVACH_FARAWAY2026_FINAL_v6.pptx"
DEVA = "Nirmala UI"  # font with Devanagari + rupee glyph coverage

prs = Presentation(SRC)
slides = list(prs.slides)

def runs(shape):
    if not shape.has_text_frame:
        return []
    return [r for p in shape.text_frame.paragraphs for r in p.runs]

def remove(shape):
    el = shape._element
    el.getparent().remove(el)

# ---------------------------------------------------------------- FIX 1: currency
print("\n===== FIX 1 — CURRENCY Rs. -> ₹ =====")
for sidx in (1, 9):  # slides 2, 10
    for sh in slides[sidx].shapes:
        for r in runs(sh):
            if "Rs." in r.text:
                before = r.text
                r.text = r.text.replace("Rs.", "₹")
                print(f"  slide {sidx+1}: '{before[:50]}' -> '{r.text[:50]}'")

# ---------------------------------------------------------------- FIX 2: arrows
print("\n===== FIX 2 — ARROWS =====")
# 2a: arrow-only '->' connector boxes on slides 3 & 6 -> proper single-line '→'
fixed_boxes = 0
for sidx in (2, 5):  # slides 3, 6
    for sh in slides[sidx].shapes:
        if sh.has_text_frame and sh.text_frame.text.strip() == "->":
            sh.text_frame.word_wrap = False
            r = runs(sh)[0]
            r.text = "→"
            r.font.name = DEVA
            fixed_boxes += 1
print(f"  connector boxes converted to '→': {fixed_boxes}")
# 2b: inline '->' in body text on slides 5, 7, 10
for sidx in (4, 6, 9):  # slides 5, 7, 10
    for sh in slides[sidx].shapes:
        for r in runs(sh):
            if "->" in r.text:
                before = r.text
                r.text = r.text.replace("->", "→")
                print(f"  slide {sidx+1} inline: '{before[:48]}' -> '{r.text[:48]}'")

# ---------------------------------------------------------------- FIX 3: slide 5 overlap
print("\n===== FIX 3 — SLIDE 5 row compression (kill badge overlap) =====")
OLD0, OLDSP = 1.92, 0.90
NEW0, NEWSP = 1.80, 0.84
moved = []
for sh in slides[4].shapes:
    t = Emu(sh.top).inches
    if 1.90 <= t < 6.90:  # the 6 tech rows (exclude title/subtitle/badges/accent)
        i = round((t - OLD0) / OLDSP)
        if 0 <= i <= 5:
            new_t = t + (NEW0 + i * NEWSP) - (OLD0 + i * OLDSP)
            moved.append((sh, t, new_t))
row_bottoms = []
for sh, old_t, new_t in moved:
    sh.top = Inches(new_t)
    row_bottoms.append(new_t + Emu(sh.height).inches)
badge_top = min(Emu(sh.top).inches for sh in slides[4].shapes if 6.90 <= Emu(sh.top).inches < 7.40)
print(f"  rows moved: {len(moved)}  | lowest row bottom now: {max(row_bottoms):.2f}\"  | badges start: {badge_top:.2f}\"")
print(f"  overlap eliminated: {max(row_bottoms) < badge_top}")

# ---------------------------------------------------------------- FIX 4: slide 6 duplicate
print("\n===== FIX 4 — SLIDE 6 delete duplicate Farmer card =====")
before_n = len(slides[5].shapes._spTree.findall('.//{http://schemas.openxmlformats.org/drawingml/2006/main}sp'))
to_del = []
for sh in slides[5].shapes:
    if sh.shape_id in (40, 41, 42, 43):  # arrow + small "Farmer Phone / Feature phone no app"
        to_del.append(sh)
for sh in to_del:
    txt = sh.text_frame.text[:30] if sh.has_text_frame else "(arrow)"
    print(f"  deleting id={sh.shape_id} '{txt}'")
    remove(sh)
print(f"  shapes removed: {len(to_del)}")

# ---------------------------------------------------------------- FIX 5: slide 9 bilingual Hindi
print("\n===== FIX 5 — SLIDE 9 bilingual notification =====")
PAIRS = [
    ("KAVACH सुरक्षा संदेश",            "KAVACH Safety Message"),
    ("तारीख: आज",                       "Date: Today"),
    ("आज का स्कोर: 95 / 100",           "Today's Score: 95 / 100"),
    ("GPS: सामान्य  |  ग्रिड: सुरक्षित",  "GPS: Normal  |  Grid: Safe"),
    ("कोई सौर तूफान नहीं",               "No solar storm"),
    ("KAVACH सक्रिय है।",                "KAVACH is active."),
]
from pptx.dml.color import RGBColor
GREEN = RGBColor(0x00, 0xff, 0x88)
MUTED = RGBColor(0x68, 0x72, 0xa8)
for sh in slides[8].shapes:
    if sh.shape_id == 47:
        print(f"  before: {repr(sh.text_frame.text)}")
        tf = sh.text_frame
        tf.clear()
        first = True
        for hi, en in PAIRS:
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            r = p.add_run(); r.text = hi
            r.font.size = Pt(11); r.font.bold = True
            r.font.color.rgb = GREEN; r.font.name = DEVA
            pe = tf.add_paragraph()
            re_ = pe.add_run(); re_.text = en
            re_.font.size = Pt(8); re_.font.italic = True
            re_.font.color.rgb = MUTED; re_.font.name = "Consolas"
            first = False
        print(f"  after : {len(PAIRS)} bilingual line-pairs")

# slide 1 कवच -> Nirmala UI
for sh in slides[0].shapes:
    for r in runs(sh):
        if "कवच" in r.text:
            r.font.name = DEVA
            print(f"\n  slide 1 कवच run -> {DEVA}")

prs.save(DST)
print(f"\nSAVED -> {DST}")
