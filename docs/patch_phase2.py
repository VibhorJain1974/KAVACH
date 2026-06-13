# Phase 2 fix: slide 9 card text overflow. v6 -> v7. Self-verifying: before/after prints.
# Slides 4 & 5 verified clean (no overflow) by render inspection -> intentionally untouched
# (raising slide 5 card height would reintroduce the badge-row overlap fixed earlier).
import sys
from pptx import Presentation
from pptx.util import Emu, Inches

sys.stdout.reconfigure(encoding="utf-8")

SRC = r"E:\FARAWAY\docs\KAVACH_FARAWAY2026_FINAL_v6.pptx"
DST = r"E:\FARAWAY\docs\KAVACH_FARAWAY2026_FINAL_v7.pptx"

prs = Presentation(SRC)
s9 = list(prs.slides)[8]
by_id = {sh.shape_id: sh for sh in s9.shapes}

def inch(emu):
    return Emu(emu).inches

# Each card: (rect, icon-label, divider, heading, body)
CARDS = [
    (7, 8, 9, 10, 11),
    (12, 13, 14, 15, 16),
    (17, 18, 19, 20, 21),
    (22, 23, 24, 25, 26),
    (27, 28, 29, 30, 31),
]
# New stack: 5 cards, height 0.90, gap 0.10 -> pitch 1.00. Start 2.00 -> last bottom 6.90 (footer 7.12).
NEW_TOP = [2.00, 3.00, 4.00, 5.00, 6.00]
CARD_H = 0.90
# Within-card offsets from card top:
OFF_LABEL = 0.10
OFF_DIV = 0.14
OFF_HEAD = 0.08
OFF_BODY = 0.42   # was 0.52 -> moved up so 2-line body ends ~+0.82, leaving padding
BODY_H = 0.42     # room for 2 lines @ 12pt (2*0.20)

print("===== SLIDE 9 — card overflow fix =====")
print(f"{'card':<6}{'rect.top b/a':<18}{'rect.h b/a':<16}{'body.top b/a':<18}{'body.h b/a':<14}")
for i, (rid, lid, did, hid, bid) in enumerate(CARDS):
    rect, label, div, head, body = (by_id[x] for x in (rid, lid, did, hid, bid))
    T = NEW_TOP[i]
    b_rect_t, b_rect_h = inch(rect.top), inch(rect.height)
    b_body_t, b_body_h = inch(body.top), inch(body.height)

    rect.top = Inches(T); rect.height = Inches(CARD_H)
    label.top = Inches(T + OFF_LABEL)
    div.top = Inches(T + OFF_DIV)
    head.top = Inches(T + OFF_HEAD)
    body.top = Inches(T + OFF_BODY); body.height = Inches(BODY_H)
    body.text_frame.word_wrap = True

    print(f"{i+1:<6}{b_rect_t:>5.2f}->{T:<10.2f}{b_rect_h:>5.2f}->{CARD_H:<8.2f}"
          f"{b_body_t:>5.2f}->{T+OFF_BODY:<10.2f}{b_body_h:>5.2f}->{BODY_H:<6.2f}")

# Overlap / budget check
bottoms = [NEW_TOP[i] + CARD_H for i in range(5)]
tops = NEW_TOP
gaps_ok = all(tops[i+1] - bottoms[i] >= 0.05 for i in range(4))
body_fits = (OFF_BODY + 2 * 0.20) <= (CARD_H - 0.05)   # 2 body lines + bottom padding
print(f"\n  last card bottom: {bottoms[-1]:.2f}\"  (footer at 7.12, accent bar 7.46)")
print(f"  inter-card gaps >= 0.05: {gaps_ok}  | 2-line body fits in card: {body_fits}")

prs.save(DST)
print(f"\nSAVED -> {DST}")
