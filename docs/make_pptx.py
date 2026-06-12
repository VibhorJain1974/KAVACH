# KAVACH_FARAWAY2026_FINAL.pptx  v3
# 12 slides | dark space | real screenshots | rounded cards | gradients | shadows
import sys, os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.oxml.ns import qn
from lxml import etree

sys.stdout.reconfigure(encoding='utf-8') if hasattr(sys.stdout, 'reconfigure') else None

# ── Palette ──────────────────────────────────────────────────────────
VOID    = RGBColor(0x03, 0x08, 0x12)
CARD    = RGBColor(0x06, 0x12, 0x22)
CARD2   = RGBColor(0x08, 0x18, 0x30)
CARD3   = RGBColor(0x02, 0x0c, 0x1c)
ACCENT  = RGBColor(0x00, 0xd4, 0xff)
GREEN   = RGBColor(0x00, 0xff, 0x88)
RED     = RGBColor(0xff, 0x2d, 0x4a)
YELLOW  = RGBColor(0xff, 0xd2, 0x3f)
ORANGE  = RGBColor(0xff, 0x6b, 0x35)
WHITE   = RGBColor(0xe2, 0xe2, 0xff)
MUTED   = RGBColor(0x68, 0x72, 0xa8)
BORDER  = RGBColor(0x00, 0x2a, 0x44)
ABORDER = RGBColor(0x00, 0x5c, 0x8a)
GBORDER = RGBColor(0x00, 0x55, 0x38)
RBORDER = RGBColor(0x7a, 0x00, 0x18)
YBOARD  = RGBColor(0x88, 0x6a, 0x00)

W = Inches(13.33)
H = Inches(7.5)

BASE = r"E:\FARAWAY\docs"
SCR  = os.path.join(BASE, "screenshots")

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

# ── Helpers ───────────────────────────────────────────────────────────
def _in(v):
    return Inches(v) if isinstance(v, (int, float)) else v

def new_slide():
    s = prs.slides.add_slide(prs.slide_layouts[6])
    s.background.fill.solid()
    s.background.fill.fore_color.rgb = VOID
    return s

def tb(s, l, t, w, h, text="", sz=14, bold=False, col=WHITE,
       align=PP_ALIGN.LEFT, bg=None, bc=None, italic=False, mono=False,
       wrap=True, url=None, bc_w=1.0):
    shape = s.shapes.add_textbox(_in(l), _in(t), _in(w), _in(h))
    if bg:
        shape.fill.solid()
        shape.fill.fore_color.rgb = bg
    if bc:
        shape.line.color.rgb = bc
        shape.line.width = Pt(bc_w)
    tf = shape.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.size = Pt(sz)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = col
    r.font.name = "Consolas" if mono else "Calibri Light"
    if url:
        rPr = r._r.get_or_add_rPr()
        hlinkClick = etree.SubElement(rPr, qn('a:hlinkClick'))
        rId = shape.part.relate_to(
            url,
            'http://schemas.openxmlformats.org/officeDocument/2006/relationships/hyperlink',
            is_external=True)
        hlinkClick.set(qn('r:id'), rId)
    return shape

def rect(s, l, t, w, h, bg=CARD, bc=None, bc_w=1.0):
    shape = s.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if bc:
        shape.line.color.rgb = bc
        shape.line.width = Pt(bc_w)
    else:
        shape.line.fill.background()
    return shape

def rrect(s, l, t, w, h, bg=CARD, bc=None, bc_w=1.0):
    shape = s.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        _in(l), _in(t), _in(w), _in(h))
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg
    if bc:
        shape.line.color.rgb = bc
        shape.line.width = Pt(bc_w)
    else:
        shape.line.fill.background()
    return shape

def pic(s, fname, l, t, w, h=None):
    path = os.path.join(SCR, fname)
    if not os.path.exists(path):
        ph = rrect(s, l, t, w, h or w * 0.625, bg=CARD2, bc=BORDER)
        tb(s, l + 0.2, t + (h or w * 0.625) / 2 - 0.2, w - 0.4, 0.4,
           "[screenshot not found]", sz=9, col=MUTED, align=PP_ALIGN.CENTER, mono=True)
        return ph
    try:
        img = s.shapes.add_picture(path, _in(l), _in(t), _in(w))
        if h is not None:
            img.height = _in(h)
        return img
    except Exception as e:
        print(f"Warning: {fname}: {e}")
        return None

def pic_framed(s, fname, l, t, w, h, border_col=ACCENT, bw=2.0):
    """Embed screenshot then draw transparent border on top."""
    img = pic(s, fname, l, t, w, h)
    frame = s.shapes.add_shape(1, _in(l), _in(t), _in(w), _in(h))
    frame.fill.background()
    frame.line.color.rgb = border_col
    frame.line.width = Pt(bw)
    return img

def label(s, text):
    tb(s, 0.5, 0.32, 12.0, 0.38, text, sz=9, col=ACCENT, mono=True)

def slide_num(s, n, total=12):
    tb(s, 11.9, 7.12, 1.3, 0.34, f"{n:02d}/{total}", sz=8, col=MUTED,
       align=PP_ALIGN.RIGHT, mono=True)

def accent_bar(s, top=7.46, height=0.03, col=ACCENT):
    rect(s, 0.0, top, 13.33, height, bg=col)

def top_bar(s, col=ACCENT):
    rect(s, 0.0, 0.0, 13.33, 0.04, bg=col)

def add_gradient_fill(shape, c1, c2, angle_deg=180):
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    for child in list(spPr):
        tag = child.tag.split('}')[-1] if '}' in child.tag else child.tag
        if tag in ('solidFill', 'gradFill', 'noFill', 'pattFill', 'blipFill'):
            spPr.remove(child)
    ang = int(angle_deg * 60000)
    xml = (
        '<a:gradFill xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        '<a:gsLst>'
        f'<a:gs pos="0"><a:srgbClr val="{c1}"/></a:gs>'
        f'<a:gs pos="100000"><a:srgbClr val="{c2}"/></a:gs>'
        '</a:gsLst>'
        f'<a:lin ang="{ang}" scaled="0"/>'
        '</a:gradFill>'
    )
    ln = spPr.find(qn('a:ln'))
    if ln is not None:
        ln.addprevious(etree.fromstring(xml))
    else:
        spPr.append(etree.fromstring(xml))

def add_shadow(shape, alpha=28000, blur=63500, dist=30000):
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is None:
        return
    if spPr.find(qn('a:effectLst')) is not None:
        return
    xml = (
        '<a:effectLst xmlns:a="http://schemas.openxmlformats.org/drawingml/2006/main">'
        f'<a:outerShdw blurRad="{blur}" dist="{dist}" dir="5400000" algn="ctr" rotWithShape="0">'
        f'<a:srgbClr val="000000"><a:alpha val="{alpha}"/></a:srgbClr>'
        '</a:outerShdw>'
        '</a:effectLst>'
    )
    spPr.append(etree.fromstring(xml))

def stat_card(s, l, t, w, h, val, lbl, col=ACCENT):
    card = rrect(s, l, t, w, h, bg=CARD2, bc=BORDER)
    add_shadow(card)
    tb(s, l + 0.12, t + 0.1, w - 0.22, h * 0.62,
       val, sz=38, bold=True, col=col, mono=True, align=PP_ALIGN.CENTER)
    tb(s, l + 0.1, t + h * 0.65, w - 0.18, h * 0.35,
       lbl, sz=9, col=MUTED, align=PP_ALIGN.CENTER)

def section_title(s, title, sub=None):
    tb(s, 0.5, 0.5, 8.5, 0.82, title, sz=46, bold=True, col=WHITE)
    if sub:
        tb(s, 0.5, 1.38, 10.5, 0.45, sub, sz=13, col=MUTED, italic=True)

def feature_row(s, l, t, w, h, code, title, desc, col=ACCENT):
    card = rrect(s, l, t, w, h, bg=CARD2, bc=BORDER)
    add_shadow(card)
    tb(s, l + 0.15, t + 0.1, 1.05, h - 0.18,
       code, sz=15, bold=True, col=col, mono=True, align=PP_ALIGN.CENTER)
    rect(s, l + 1.28, t + 0.14, 0.02, h - 0.26, bg=BORDER)
    tb(s, l + 1.38, t + 0.08, w - 1.55, 0.44, title, sz=13, bold=True, col=col)
    tb(s, l + 1.38, t + 0.52, w - 1.55, h - 0.62, desc, sz=12, col=WHITE)


# ====================================================================
# SLIDE 1 — TITLE
# ====================================================================
s = new_slide()
slide_num(s, 1)
top_bar(s)

# Background gradient overlay
bg = rect(s, 0.0, 0.04, 13.33, 7.42, bg=VOID)
add_gradient_fill(bg, "030812", "041628", 145)

# Eyebrow strip
strip = rect(s, 0.0, 0.04, 13.33, 0.46, bg=CARD3)
tb(s, 0.55, 0.1, 9.5, 0.35,
   "FAR AWAY 2026  |  SPACE & AEROSPACE  |  TEAM 404_SHINOBI",
   sz=9, col=ACCENT, mono=True)
tb(s, 11.0, 0.1, 2.15, 0.35, "14 JUN 2026",
   sz=9, col=MUTED, mono=True, align=PP_ALIGN.RIGHT)

# Hero title
tb(s, 0.48, 0.72, 8.0, 2.2, "KAVACH", sz=112, bold=True, col=WHITE)

# Kavach line under title
accent_line = rect(s, 0.48, 2.88, 7.5, 0.05, bg=ACCENT)
add_gradient_fill(accent_line, "00d4ff", "004466", 0)

# Hindi + tagline
tb(s, 0.52, 3.05, 9.5, 0.52,
   "कवच  —  Autonomous Space Weather Shield for India",
   sz=17, col=ACCENT, italic=True)
tb(s, 0.52, 3.65, 9.8, 0.58,
   '"It called the farmer before the lights went out."',
   sz=18, col=MUTED, italic=True)

# Three stat pills
for i, (val, lbl, col) in enumerate([
    ("28",     "DISCOMs warned",      ACCENT),
    ("12,400", "Farmers per storm",   GREEN),
    ("0",      "Human triggers",      YELLOW),
]):
    x = 0.5 + i * 4.28
    stat_card(s, x, 4.5, 3.95, 1.55, val, lbl, col)

# Live URL row
rrect(s, 0.5, 6.35, 5.85, 0.38, bg=CARD3, bc=ABORDER)
tb(s, 0.66, 6.4, 5.65, 0.28,
   "LIVE: frontend-rust-xi-79.vercel.app", sz=9, col=ACCENT, mono=True)
rrect(s, 6.52, 6.35, 6.32, 0.38, bg=CARD3, bc=BORDER)
tb(s, 6.68, 6.4, 6.1, 0.28,
   "REPO: github.com/VibhorJain1974/kavach-faraway-2026", sz=9, col=MUTED, mono=True)

accent_bar(s)


# ====================================================================
# SLIDE 2 — PROBLEM
# ====================================================================
s = new_slide()
label(s, "PROBLEM")
slide_num(s, 2)
top_bar(s, RED)

# Hero statement
tb(s, 0.5, 0.5, 12.0, 1.1, "Zero Warning.", sz=58, bold=True, col=RED)

# Date + severity stamp
stamp = rrect(s, 0.5, 1.65, 12.5, 0.55, bg=RBORDER)
add_gradient_fill(stamp, "4a000e", "1a0008", 0)
tb(s, 0.68, 1.74, 12.0, 0.35,
   "May 10-11, 2024  |  Kp = 9.0  |  G5 EXTREME  |  India's grid nearly failed  |  0 warnings issued",
   sz=12, bold=True, col=RED, mono=True)

# Three problem cards
for i, (icon, title, headline, detail, col, bc) in enumerate([
    ("!!", "No Indian System",   "28 DISCOMs",
     "Zero automated alerts\nacross all power utilities",       RED,    RBORDER),
    ("X",  "Wrong Language",     "English-only",
     "Internet-only western tools\nuseless at India's last mile", YELLOW, YBOARD),
    ("?",  "Feature Phone Gap",  "300M farmers",
     "No smartphone, no internet\nno way to receive any warning", ORANGE, RGBColor(0x66,0x33,0x00)),
]):
    x = 0.5 + i * 4.28
    card = rrect(s, x, 2.45, 3.92, 3.55, bg=CARD2, bc=bc, bc_w=1.5)
    add_shadow(card, alpha=35000)
    tb(s, x + 0.2, 2.6, 0.7, 0.7, icon, sz=28, bold=True, col=col, mono=True)
    tb(s, x + 0.2, 3.3, 3.5, 0.45, title, sz=14, bold=True, col=col)
    tb(s, x + 0.2, 3.78, 3.5, 0.44, headline, sz=22, bold=True, col=WHITE)
    tb(s, x + 0.2, 4.25, 3.5, 1.55, detail, sz=14, col=MUTED)

# Bottom impact line
impact = rrect(s, 0.5, 6.25, 12.45, 0.72, bg=CARD3, bc=RBORDER)
tb(s, 0.65, 6.35, 12.0, 0.55,
   "Rs.3,000 Cr transformer assets at risk every G5 event  |  4.5 hours of actionable warning window lost forever",
   sz=12, col=RED, italic=True, align=PP_ALIGN.CENTER)

accent_bar(s, col=RED)


# ====================================================================
# SLIDE 3 — SOLUTION
# ====================================================================
s = new_slide()
label(s, "SOLUTION")
slide_num(s, 3)
top_bar(s, ACCENT)

tb(s, 0.5, 0.5, 10.0, 0.92, "KAVACH", sz=58, bold=True, col=ACCENT)
tb(s, 0.5, 1.5, 12.0, 0.48,
   "Autonomous  |  Real-time  |  Hindi voice  |  No human trigger",
   sz=14, col=MUTED, mono=True)

# Flow diagram
nodes = [
    ("NASA\nDONKI",   "15 min",  WHITE,  BORDER),
    ("NOAA\nSWPC",    "5 min",   WHITE,  BORDER),
    ("KAVACH\nBrain", "FastAPI", ACCENT, ACCENT),
    ("28\nDISCOMs",   "warned",  YELLOW, YBOARD),
    ("Farmer\nPhone", "no app",  GREEN,  GBORDER),
]
node_w, gap, flow_y = 2.05, 0.35, 2.2
flow_start = 0.42
for i, (txt, sub, col, bc) in enumerate(nodes):
    x = flow_start + i * (node_w + gap)
    card = rrect(s, x, flow_y, node_w, 1.38, bg=CARD2 if i != 2 else CARD3, bc=bc)
    if i == 2:
        add_gradient_fill(card, "001428", "002244", 180)
    add_shadow(card)
    tb(s, x + 0.1, flow_y + 0.12, node_w - 0.18, 0.72,
       txt, sz=14, bold=(i == 2), col=col, align=PP_ALIGN.CENTER)
    tb(s, x + 0.1, flow_y + 0.85, node_w - 0.18, 0.35,
       sub, sz=9, col=MUTED if i != 2 else ACCENT, mono=True, align=PP_ALIGN.CENTER)
    if i < 4:
        tb(s, x + node_w + 0.04, flow_y + 0.55, 0.33, 0.38,
           "->", sz=15, col=ACCENT, bold=True, mono=True)

# Three action blocks
for i, (num, act, detail, col) in enumerate([
    ("01", "MONITOR",  "NASA + NOAA\nevery 5–15 min\nautonomously", ACCENT),
    ("02", "MAP RISK", "28 zones\nnorth-first\nreal DISCOM data", YELLOW),
    ("03", "CALL",     "Hindi TTS\nfeature phone\n0 internet needed", GREEN),
]):
    x = 0.5 + i * 4.28
    card = rrect(s, x, 3.85, 3.92, 2.55, bg=CARD, bc=BORDER)
    add_shadow(card)
    tb(s, x + 0.18, 3.97, 0.72, 0.62, num, sz=20, bold=True, col=col, mono=True)
    rect(s, x + 0.18, 4.62, 3.55, 0.022, bg=BORDER)
    tb(s, x + 0.18, 4.72, 3.5, 0.44, act, sz=13, bold=True, col=col, mono=True)
    tb(s, x + 0.18, 5.18, 3.5, 1.1, detail, sz=14, col=WHITE)

# Hindi message
rrect(s, 0.5, 6.56, 12.45, 0.52, bg=RGBColor(0x01, 0x12, 0x08), bc=GBORDER)
tb(s, 0.68, 6.62, 12.0, 0.38,
   '"Namaste. Main KAVACH bol raha hoon. Ek bada solar toofan aa raha hai. Bijli band kar dijiye."',
   sz=11, col=GREEN, mono=True, italic=True)

accent_bar(s)


# ====================================================================
# SLIDE 4 — KEY FEATURES  (with real app screenshot)
# ====================================================================
s = new_slide()
label(s, "KEY FEATURES")
slide_num(s, 4)
top_bar(s)

tb(s, 0.5, 0.5, 6.5, 0.82, "Four layers.", sz=46, bold=True, col=WHITE)
tb(s, 0.5, 1.38, 6.0, 0.38,
   "Realtime  +  Predictive  +  Historical  +  Daily",
   sz=12, col=MUTED, mono=True)

# Left: 4 feature rows
feats = [
    ("MAP", "Command Center",
     "Live India map  |  28 DISCOM zones  |  Kp gauge  |  alert feed", ACCENT),
    ("AUR", "Aurora Predictor",
     "6 Indian locations  |  Kp-based visibility  |  Hindi alerts",    GREEN),
    ("SHD", "Daily Shield",
     "Score 0-100  |  6:30 AM IST  |  GPS + grid + aurora status",     YELLOW),
    ("MEM", "Storm Memory",
     "4 storms since 2003  |  847 alerts fired  |  counterfactual",    ORANGE),
]
for i, (code, title, desc, col) in enumerate(feats):
    y = 1.95 + i * 1.32
    feature_row(s, 0.45, y, 6.15, 1.18, code, title, desc, col)

# Right: Real app screenshot
pic_framed(s, "01_home.png", 6.82, 0.5, 6.3, 3.9375, border_col=ACCENT, bw=1.5)
# Demo in-progress screenshot below
pic_framed(s, "06_demo_red.png", 6.82, 4.55, 6.3, 2.82, border_col=RED, bw=1.5)
tb(s, 6.9, 4.62, 2.5, 0.3, "[ MAY 2024 REPLAY ]",
   sz=8, bold=True, col=RED, mono=True, bg=RGBColor(0x3a, 0x00, 0x10))

accent_bar(s)


# ====================================================================
# SLIDE 5 — TECH STACK
# ====================================================================
s = new_slide()
label(s, "TECH STACK")
slide_num(s, 5)
top_bar(s)

section_title(s, "Real APIs. Zero mocks.")
tb(s, 0.5, 1.38, 12.0, 0.42,
   "Every data source is live.  Every call is real Twilio.  Every map pin is a real DISCOM.",
   sz=12, col=MUTED, italic=True)

rows = [
    ("FRONTEND",  "Next.js 14",
     "TypeScript  +  Tailwind CSS  +  Mapbox GL JS  +  Three.js",           ACCENT),
    ("BACKEND",   "FastAPI 0.110",
     "APScheduler  +  Python 3.13  +  httpx  +  Supabase client",           ACCENT),
    ("DATABASE",  "Supabase",
     "PostgreSQL  +  realtime  |  28 DISCOMs  +  alert log  +  v2 tables",  GREEN),
    ("VOICE",     "Twilio + Polly",
     "Amazon Polly Aditi (hi-IN)  |  TTS -> MP3 fallback  |  multi-number parallel calls",    GREEN),
    ("DATA",      "NASA DONKI",
     "NOAA SWPC Kp 5-min  +  NASA CME/GST 15-min  +  local storm archive",  YELLOW),
    ("DEPLOY",    "Vercel + Railway",
     "frontend-rust-xi-79.vercel.app  +  powerful-respect-production.up.railway.app", MUTED),
]
for i, (layer, ver, detail, col) in enumerate(rows):
    y = 1.92 + i * 0.9
    # Layer pill
    pill = rrect(s, 0.45, y, 1.92, 0.76, bg=CARD2, bc=BORDER)
    tb(s, 0.55, y + 0.16, 1.74, 0.42, layer, sz=11, bold=True, col=col, mono=True, align=PP_ALIGN.CENTER)
    # Version chip
    ver_chip = rrect(s, 2.48, y, 1.72, 0.76, bg=CARD3, bc=BORDER)
    tb(s, 2.58, y + 0.16, 1.56, 0.42, ver, sz=11, bold=True, col=col, mono=True, align=PP_ALIGN.CENTER)
    # Detail
    detail_r = rrect(s, 4.32, y, 8.6, 0.76, bg=CARD, bc=BORDER)
    tb(s, 4.46, y + 0.18, 8.34, 0.42, detail, sz=12, col=WHITE, mono=True)

# Bottom proof badges — inline only, no stat_card (avoids sz=38 overflow)
for i, (val, lbl, col) in enumerate([
    ("LIVE",   "Vercel deploy",    GREEN),
    ("LIVE",   "Railway backend",  GREEN),
    ("LIVE",   "Twilio tested",    GREEN),
    ("6 DAYS", "Concept to ship",  ACCENT),
]):
    x = 0.45 + i * 3.22
    rrect(s, x, 6.94, 3.0, 0.44, bg=CARD3, bc=GBORDER if col == GREEN else ABORDER)
    tb(s, x + 0.14, 7.04, 0.85, 0.26, val, sz=10, bold=True, col=col, mono=True)
    tb(s, x + 1.06, 7.04, 1.88, 0.26, lbl, sz=10, col=MUTED)

accent_bar(s)


# ====================================================================
# SLIDE 6 — ARCHITECTURE
# ====================================================================
s = new_slide()
label(s, "ARCHITECTURE")
slide_num(s, 6)
top_bar(s)

section_title(s, "Autonomous loop.")
tb(s, 0.5, 1.38, 12.0, 0.42,
   "No human in the trigger path.  Polls every 5-15 min.  Fires alerts in parallel.",
   sz=12, col=MUTED, italic=True)

# Left: data sources
for i, (src, freq, col, bc) in enumerate([
    ("NASA DONKI", "CME + GST\n15 min",   WHITE,  ABORDER),
    ("NOAA SWPC",  "Kp-index\n5 min",     WHITE,  ABORDER),
    ("ISRO VEDAS", "Phase 2\nionosphere", MUTED,  BORDER),
]):
    y = 1.92 + i * 1.72
    card = rrect(s, 0.38, y, 2.35, 1.42, bg=CARD2, bc=bc)
    if col != MUTED:
        add_shadow(card)
    tb(s, 0.5, y + 0.14, 2.1, 0.52, src, sz=13, bold=True, col=col)
    tb(s, 0.5, y + 0.68, 2.1, 0.55, freq, sz=10, col=MUTED, mono=True)
    if col != MUTED:
        tb(s, 2.82, y + 0.52, 0.38, 0.42, "->", sz=15, col=ACCENT, bold=True, mono=True)

# KAVACH brain - center
brain = rrect(s, 3.28, 1.78, 3.35, 5.0, bg=CARD3, bc=ACCENT, bc_w=1.5)
add_gradient_fill(brain, "001428", "020c1c", 180)
add_shadow(brain, alpha=40000)
rect(s, 3.28, 1.78, 3.35, 0.06, bg=ACCENT)  # accent top line
tb(s, 3.4, 1.9, 3.12, 0.62, "KAVACH", sz=28, bold=True, col=ACCENT, mono=True)
tb(s, 3.4, 2.52, 3.12, 0.38, "FastAPI + APScheduler", sz=10, col=MUTED, mono=True)
rect(s, 3.4, 2.98, 3.1, 0.022, bg=BORDER)
for j, comp in enumerate(["Storm Classifier", "DISCOM Mapper (x28)",
                           "Aurora Predictor", "Daily Shield", "Alert Generator"]):
    tb(s, 3.42, 3.08 + j * 0.65, 3.1, 0.48, "  " + comp, sz=11, col=WHITE)

# Right: outputs
outputs = [
    ("Supabase",     "Alert log + history",   ACCENT,  ABORDER),
    ("Twilio",       "Hindi voice call",       GREEN,   GBORDER),
    ("Next.js Map",  "28 zones realtime",      YELLOW,  YBOARD),
    ("Farmer Phone", "Feature phone  no app",  GREEN,   GBORDER),
]
for i, (out, sub, col, bc) in enumerate(outputs):
    y = 1.92 + i * 1.42
    tb(s, 6.82, y + 0.42, 0.38, 0.42, "->", sz=14, col=col, bold=True, mono=True)
    card = rrect(s, 7.28, y, 2.55, 1.18, bg=CARD, bc=bc)
    add_shadow(card)
    tb(s, 7.42, y + 0.12, 2.3, 0.44, out, sz=13, bold=True, col=col)
    tb(s, 7.42, y + 0.58, 2.3, 0.42, sub, sz=10, col=MUTED)

# Final endpoint
farmer = rrect(s, 10.08, 2.72, 2.92, 2.1, bg=RGBColor(0x01, 0x14, 0x0a), bc=GREEN)
add_shadow(farmer)
tb(s, 10.2, 2.85, 2.7, 0.52, "Farmer's Phone", sz=14, bold=True, col=GREEN)
tb(s, 10.2, 3.4, 2.7, 1.0,
   "No app needed\nNo internet\nJust a call", sz=12, col=MUTED)
tb(s, 9.72, 3.55, 0.38, 0.42, "->", sz=14, col=GREEN, bold=True, mono=True)

accent_bar(s)


# ====================================================================
# SLIDE 7 — DEMO  (big screenshot + timeline)
# ====================================================================
s = new_slide()
label(s, "DEMO")
slide_num(s, 7)
top_bar(s)

# Red LIVE badge
live_badge = rrect(s, 0.5, 0.5, 2.65, 0.75, bg=RED)
add_gradient_fill(live_badge, "cc0022", "880011", 180)
add_shadow(live_badge, alpha=50000)
tb(s, 0.6, 0.56, 2.45, 0.58, "[ LIVE DEMO ]", sz=15, bold=True,
   col=WHITE, align=PP_ALIGN.CENTER, mono=True)

tb(s, 3.42, 0.5, 8.0, 0.82, "22.6 seconds.", sz=46, bold=True, col=WHITE)
tb(s, 3.42, 1.38, 6.0, 0.38,
   "May 10, 2024  |  Kp = 9.0  |  G5 EXTREME",
   sz=12, col=MUTED, mono=True)

# Left: 6-step timeline
steps = [
    ("T+0s",  "NOMINAL",    "Monitoring active",                WHITE,  BORDER),
    ("T+3s",  "CME",        "X-class flare detected",           YELLOW, YBOARD),
    ("T+5s",  "Kp = 9.0",   "EXTREME  G5  confirmed",           RED,    RBORDER),
    ("T+14s", "MAP RED",    "All 28 zones -> red",              RED,    RBORDER),
    ("T+17s", "ALERTS",     "28 DISCOMs  +  12,400 farmers",    ORANGE, RGBColor(0x66,0x33,0x00)),
    ("T+20s", "CALL LIVE",  "Hindi rings on demo phone",         GREEN,  GBORDER),
]
for i, (t, status, desc, col, bc) in enumerate(steps):
    y = 1.92 + i * 0.87
    # Time box
    time_r = rrect(s, 0.45, y, 1.08, 0.72, bg=CARD2, bc=BORDER)
    tb(s, 0.52, y + 0.18, 0.95, 0.36, t, sz=10, col=ACCENT, mono=True, align=PP_ALIGN.CENTER)
    # Status box
    status_r = rrect(s, 1.65, y, 2.12, 0.72, bg=CARD2, bc=bc)
    tb(s, 1.76, y + 0.18, 1.95, 0.36, status, sz=11, bold=True, col=col, mono=True)
    # Description
    desc_r = rrect(s, 3.9, y, 3.05, 0.72, bg=CARD, bc=BORDER)
    tb(s, 4.04, y + 0.18, 2.82, 0.36, desc, sz=12, col=WHITE)
    # Connector line (not for last)
    if i < 5:
        rect(s, 0.98, y + 0.72, 0.02, 0.15, bg=MUTED)

# Right: Large demo screenshot - the money shot
pic_framed(s, "06_demo_red.png", 7.22, 1.78, 5.85, 3.66, border_col=RED, bw=2.0)

# Caption under screenshot
tb(s, 7.22, 5.52, 5.85, 0.32,
   "India map: all 28 DISCOMs in RED  |  Alerts firing in parallel",
   sz=9, col=RED, mono=True, italic=True)

# Stats below timeline — simple badges, no stat_card
for j, (val, lbl, col) in enumerate([
    ("22.6s", "demo runtime",     YELLOW),
    ("4.5h",  "warning window",   ACCENT),
    ("LIVE",  "Vercel + Railway",  GREEN),
]):
    x = 0.45 + j * 2.3
    rrect(s, x, 7.0, 2.15, 0.42, bg=CARD2, bc=BORDER)
    tb(s, x + 0.12, 7.08, 0.82, 0.26, val, sz=11, bold=True, col=col, mono=True)
    tb(s, x + 0.98, 7.08, 1.1, 0.26, lbl, sz=9, col=MUTED)

accent_bar(s, col=RED)


# ====================================================================
# SLIDE 8 — STORM MEMORY
# ====================================================================
s = new_slide()
label(s, "STORM MEMORY")
slide_num(s, 8)
top_bar(s)

section_title(s, "21 years of storms.")
tb(s, 0.5, 1.38, 12.0, 0.42,
   "KAVACH applied retroactively to every major Indian geomagnetic event since 2003.",
   sz=12, col=MUTED, italic=True)

storms = [
    ("2003", "Halloween Superstorm",  "Kp 9.0+", "8h avg",  "294", ORANGE),
    ("2015", "St. Patrick's Day",     "Kp 8.0",  "6h avg",  "216", ACCENT),
    ("2017", "September X9.3 Flare",  "Kp 8.0",  "5h avg",  "192", YELLOW),
    ("2024", "May G5  [SOURCE]",       "Kp 9.0",  "4.5h",   "145", RED),
]
for i, (yr, name, kp, warn, alerts, col) in enumerate(storms):
    y = 2.0 + i * 1.0
    # Year
    yr_r = rrect(s, 0.38, y, 0.95, 0.82, bg=CARD2, bc=col)
    tb(s, 0.4, y + 0.2, 0.92, 0.42, yr, sz=14, bold=True, col=col,
       align=PP_ALIGN.CENTER, mono=True)
    # Name
    name_r = rrect(s, 1.45, y, 5.95, 0.82, bg=CARD, bc=BORDER)
    tb(s, 1.6, y + 0.2, 5.7, 0.42, name, sz=14, col=col if yr == "2024" else WHITE)
    # Kp
    kp_r = rrect(s, 7.52, y, 1.6, 0.82, bg=CARD2, bc=BORDER)
    tb(s, 7.62, y + 0.2, 1.44, 0.42, kp, sz=12, col=MUTED, mono=True, align=PP_ALIGN.CENTER)
    # Warning window
    warn_r = rrect(s, 9.24, y, 2.1, 0.82, bg=RGBColor(0x01, 0x14, 0x0a), bc=GREEN)
    tb(s, 9.34, y + 0.18, 1.9, 0.44, warn, sz=14, bold=True, col=GREEN, mono=True, align=PP_ALIGN.CENTER)
    # Alerts fired
    al_r = rrect(s, 11.46, y, 1.5, 0.82, bg=CARD2, bc=BORDER)
    tb(s, 11.56, y + 0.2, 1.3, 0.42, alerts + " alerts", sz=11, col=col, mono=True, align=PP_ALIGN.CENTER)

# Column headers
for x, w2, txt in [(0.38, 0.95, "YEAR"), (1.45, 5.95, "STORM"),
                   (7.52, 1.6, "Kp"), (9.24, 2.1, "WARNING"),
                   (11.46, 1.5, "ALERTS")]:
    tb(s, x, 1.72, w2, 0.28, txt, sz=8, col=MUTED, mono=True, align=PP_ALIGN.CENTER)

# Three big proof stats
for i, (val, lbl, col) in enumerate([
    ("847",    "total alerts fired",          RED),
    ("12,400", "subscribers per storm",        YELLOW),
    ("6h",     "average warning window",       GREEN),
]):
    x = 0.38 + i * 4.32
    stat_card(s, x, 6.28, 4.1, 1.0, val, lbl, col)

accent_bar(s)


# ====================================================================
# SLIDE 9 — DAILY SHIELD
# ====================================================================
s = new_slide()
label(s, "DAILY SHIELD")
slide_num(s, 9)
top_bar(s)

section_title(s, "365 days a year.", "Space Weather Score  |  6:30 AM IST  |  Hindi + English")

# Left: features
for i, (code, title, desc, col) in enumerate([
    ("0-100", "Space Weather Score",  "India's space weather AQI\nGreen 80+ = normal ops",       ACCENT),
    ("GPS",   "GPS Accuracy Forecast","Precision ag tractors\nFishing boat GPS alert",            GREEN),
    ("GRID",  "DISCOM Risk Zones",    "Per-zone transformer risk\nOperator dashboard feed",       YELLOW),
    ("AUR",   "Aurora Visibility",    "6 Himalayan locations\nKp-based prediction nightly",       GREEN),
    ("LANG",  "Bilingual Brief",      "Hindi for farmers\nEnglish for DISCOM operators",         WHITE),
]):
    y = 2.0 + i * 0.98
    feature_row(s, 0.42, y, 6.12, 0.86, code, title, desc, col)

# Right: phone mockup card
phone_l, phone_t, phone_w, phone_h = 6.82, 0.62, 6.3, 6.62
phone = rrect(s, phone_l, phone_t, phone_w, phone_h, bg=CARD3, bc=ACCENT, bc_w=1.5)
add_gradient_fill(phone, "020c1c", "03101e", 180)
add_shadow(phone, alpha=45000)
rect(s, phone_l, phone_t, phone_w, 0.06, bg=ACCENT)

# App header bar
app_hdr = rect(s, phone_l, phone_t + 0.06, phone_w, 0.48, bg=CARD2)
tb(s, phone_l + 0.2, phone_t + 0.12, 3.5, 0.36,
   "KAVACH  [ 06:30 AM IST ]", sz=11, col=ACCENT, mono=True, bold=True)
tb(s, phone_l + 4.8, phone_t + 0.12, 1.3, 0.36,
   "NOMINAL", sz=10, col=GREEN, mono=True, align=PP_ALIGN.RIGHT)

# Score gauge
score_bg = rrect(s, phone_l + 0.2, phone_t + 0.68, phone_w - 0.38, 1.9, bg=RGBColor(0x01,0x10,0x06), bc=GREEN)
tb(s, phone_l + 0.28, phone_t + 0.78, phone_w - 0.55, 1.2,
   "95", sz=80, bold=True, col=GREEN, mono=True, align=PP_ALIGN.CENTER)
tb(s, phone_l + 0.28, phone_t + 1.9, phone_w - 0.55, 0.42,
   "/ 100   SHIELD NOMINAL", sz=12, col=GREEN, mono=True, align=PP_ALIGN.CENTER)

# Status chips row
for j, (lbl, col) in enumerate([("GPS OK", GREEN), ("GRID OK", GREEN), ("AURORA LOW", MUTED)]):
    cx = phone_l + 0.22 + j * 2.02
    chip_r = rrect(s, cx, phone_t + 2.76, 1.9, 0.38, bg=CARD2, bc=GBORDER)
    tb(s, cx + 0.08, phone_t + 2.84, 1.75, 0.24, lbl, sz=9, col=col, mono=True, align=PP_ALIGN.CENTER)

# Hindi notification box
notif = rrect(s, phone_l + 0.2, phone_t + 3.3, phone_w - 0.38, 2.58, bg=RGBColor(0x02, 0x0e, 0x1a), bc=BORDER)
tb(s, phone_l + 0.32, phone_t + 3.42, phone_w - 0.6, 2.35,
   "KAVACH Suraksha Sandesh\n"
   "Tanikh: Aaj\n"
   "Aaj ka Score: 95 / 100\n"
   "GPS: Sahi  |  Grid: Surakshit\n"
   "Koi solar toofan nahi\n"
   "KAVACH active hai.",
   sz=12, col=GREEN, mono=True)

# Score badge next to phone
score_stamp = rrect(s, phone_l + 0.2, phone_t + 6.0, phone_w - 0.38, 0.42, bg=GBORDER)
tb(s, phone_l + 0.28, phone_t + 6.06, phone_w - 0.55, 0.3,
   "INDIA'S SPACE WEATHER AQI  |  DAILY AT 06:30 IST",
   sz=9, col=GREEN, mono=True, align=PP_ALIGN.CENTER)

accent_bar(s)


# ====================================================================
# SLIDE 10 — REVENUE
# ====================================================================
s = new_slide()
label(s, "REVENUE")
slide_num(s, 10)
top_bar(s, GREEN)

# 600x ROI hero
roi_hero = rrect(s, 0.45, 0.5, 12.45, 1.05, bg=RGBColor(0x01, 0x14, 0x0a), bc=GREEN, bc_w=2.0)
add_gradient_fill(roi_hero, "011408", "001c08", 0)
add_shadow(roi_hero, alpha=40000)
tb(s, 0.62, 0.58, 12.0, 0.82,
   "600x  ROI",
   sz=52, bold=True, col=GREEN, mono=True, align=PP_ALIGN.CENTER)

tb(s, 0.5, 1.65, 12.0, 0.42,
   "Rs.5L/month subscription  protects  Rs.3,000 Cr in transformer assets per DISCOM",
   sz=13, col=MUTED, italic=True, align=PP_ALIGN.CENTER)

# Tier table
tiers = [
    ("FREE",        "Citizens / Farmers / Fishermen",
     "Public alerts + Hindi voice calls",       "Rs. 0",          WHITE,   BORDER),
    ("DISCOM API",  "28 Electricity Utilities",
     "72h forecast + SCADA integration + Portal", "Rs.2-5L/month", ACCENT,  ABORDER),
    ("5 DISCOMs",   "Year 1 anchor contracts (target)",
     "Pilot -> ARR target -> Round 2 expansion", "Rs.1.8Cr/year",  GREEN,   GBORDER),
    ("PHASE 2",     "Aviation / Telecom / Insurance",
     "DGCA + BSNL + LIC APIs as risk oracles",  "TBD",            YELLOW,  YBOARD),
]
for i, (tier, who, detail, price, col, bc) in enumerate(tiers):
    y = 2.28 + i * 1.1
    tier_r = rrect(s, 0.45, y, 2.12, 0.9, bg=CARD2, bc=bc)
    tb(s, 0.55, y + 0.2, 1.95, 0.48, tier, sz=12, bold=True, col=col, mono=True, align=PP_ALIGN.CENTER)
    who_r = rrect(s, 2.7, y, 3.6, 0.9, bg=CARD, bc=BORDER)
    tb(s, 2.84, y + 0.2, 3.38, 0.48, who, sz=12, col=WHITE)
    detail_r = rrect(s, 6.42, y, 4.38, 0.9, bg=CARD, bc=BORDER)
    tb(s, 6.56, y + 0.2, 4.14, 0.48, detail, sz=11, col=MUTED)
    price_r = rrect(s, 10.92, y, 2.0, 0.9, bg=CARD2, bc=bc)
    tb(s, 11.0, y + 0.2, 1.84, 0.48, price, sz=12, bold=True, col=col, mono=True, align=PP_ALIGN.CENTER)

# Visual cost/benefit bar
tb(s, 0.5, 6.75, 12.0, 0.42,
   "Rs.3,000 Cr  [transformer asset exposure]  vs  Rs.5L  [monthly subscription]  =  600x protection",
   sz=11, col=GREEN, mono=True, italic=True, align=PP_ALIGN.CENTER)
bar_bg = rect(s, 0.45, 6.62, 12.45, 0.06, bg=BORDER)
bar_fill = rect(s, 0.45, 6.62, 12.43, 0.06, bg=GREEN)
add_gradient_fill(bar_fill, "00ff88", "004422", 0)

accent_bar(s, col=GREEN)


# ====================================================================
# SLIDE 11 — FUTURE SCOPE
# ====================================================================
s = new_slide()
label(s, "FUTURE SCOPE")
slide_num(s, 11)
top_bar(s)

tb(s, 0.5, 0.5, 12.0, 0.82, "Round 2 / Tokyo.", sz=46, bold=True, col=WHITE)

# Flagship banner
flagship = rrect(s, 0.42, 1.48, 12.48, 1.22, bg=RGBColor(0x01, 0x14, 0x0a), bc=GREEN, bc_w=1.5)
add_gradient_fill(flagship, "011408", "001c08", 135)
add_shadow(flagship)
rect(s, 0.42, 1.48, 12.48, 0.055, bg=GREEN)
tb(s, 0.6, 1.56, 2.8, 0.32, "FLAGSHIP  #1", sz=9, col=GREEN, mono=True, bold=True)
tb(s, 0.6, 1.86, 12.1, 0.68,
   "Multi-Language:  Hindi  +  Tamil  +  Telugu  +  Marathi  +  Bengali  +  Punjabi  +  Gujarati  +  Kannada",
   sz=14, bold=True, col=GREEN)

# 4 cards in 2x2
feats2 = [
    ("VEDAS",  "ISRO VEDAS Integration",  "India-specific ionospheric\nscintillation data",    ACCENT,  ABORDER),
    ("WA/SMS", "WhatsApp + SMS Alerts",   "Multi-channel delivery\nzero missed notifications", YELLOW,  YBOARD),
    ("REPLY",  "Two-Way Voice Reply",     "Farmer presses 1 to confirm\n2 for help — KAVACH escalates",  GREEN,   GBORDER),
    ("CME",    "CME Transit Model",       "4-18h precise countdown\nper-zone impact window",    ORANGE,  RGBColor(0x55,0x22,0x00)),
]
for i, (code, title, desc, col, bc) in enumerate(feats2):
    col_n = i % 2
    row_n = i // 2
    x = 0.42 + col_n * 6.28
    y = 2.92 + row_n * 1.65
    card = rrect(s, x, y, 5.98, 1.48, bg=CARD2, bc=bc)
    add_shadow(card)
    tb(s, x + 0.18, y + 0.12, 1.05, 0.55, code, sz=15, bold=True, col=col, mono=True)
    rect(s, x + 1.32, y + 0.14, 0.025, 0.55, bg=BORDER)
    tb(s, x + 1.44, y + 0.12, 4.38, 0.42, title, sz=13, bold=True, col=col)
    tb(s, x + 1.44, y + 0.58, 4.38, 0.78, desc, sz=12, col=MUTED)

# Bottom full-width card
saas = rrect(s, 0.42, 6.22, 12.48, 0.98, bg=CARD, bc=ABORDER)
add_shadow(saas)
tb(s, 0.6, 6.3, 1.2, 0.55, "SaaS", sz=16, bold=True, col=ACCENT, mono=True)
rect(s, 1.88, 6.32, 0.025, 0.55, bg=BORDER)
tb(s, 2.0, 6.3, 4.5, 0.42, "DISCOM Operator Dashboard", sz=13, bold=True, col=ACCENT)
tb(s, 2.0, 6.74, 10.7, 0.32, "Operator-facing SaaS portal  |  Custom alerting  |  SCADA integration  |  Contractual ARR path", sz=11, col=MUTED)

accent_bar(s)


# ====================================================================
# SLIDE 12 — THANK YOU
# ====================================================================
s = new_slide()
slide_num(s, 12)
top_bar(s)

bg2 = rect(s, 0.0, 0.04, 13.33, 7.42, bg=VOID)
add_gradient_fill(bg2, "030812", "041628", 145)

# Top badge
tb(s, 0.5, 0.22, 12.0, 0.38,
   "FAR AWAY 2026  |  TEAM 404_SHINOBI  |  SPACE & AEROSPACE",
   sz=9, col=ACCENT, mono=True)

# Hero
tb(s, 0.48, 0.72, 9.5, 2.1, "KAVACH", sz=112, bold=True, col=WHITE)
accent_line2 = rect(s, 0.48, 2.88, 7.5, 0.05, bg=ACCENT)
add_gradient_fill(accent_line2, "00d4ff", "003355", 0)
tb(s, 0.52, 3.05, 10.5, 0.52,
   '"It called the farmer before the lights went out."',
   sz=18, col=MUTED, italic=True)

# 4 link cards 2x2
links = [
    ("LIVE DEMO",   "frontend-rust-xi-79.vercel.app",
     "https://frontend-rust-xi-79.vercel.app",                        ACCENT,  ABORDER),
    ("GITHUB REPO", "github.com/VibhorJain1974/kavach-faraway-2026",
     "https://github.com/VibhorJain1974/kavach-faraway-2026",          ACCENT,  ABORDER),
    ("CONTACT",     "jvibhor202@gmail.com",
     "mailto:jvibhor202@gmail.com",                                    MUTED,   BORDER),
    ("TEAM",        "404_SHINOBI  |  FAR AWAY 2026",
     None,                                                              YELLOW,  YBOARD),
]
for i, (lbl, val, url, col, bc) in enumerate(links):
    col_n = i % 2
    row_n = i // 2
    x = 0.5 + col_n * 6.5
    y = 3.82 + row_n * 1.38
    card = rrect(s, x, y, 6.12, 1.22, bg=CARD2, bc=bc)
    add_shadow(card)
    tb(s, x + 0.18, y + 0.1, 5.8, 0.34, lbl, sz=8, col=MUTED, mono=True)
    tb(s, x + 0.18, y + 0.48, 5.8, 0.56, val, sz=14, col=col, url=url)

# Small proof screenshot
pic_framed(s, "07_demo_complete.png", 9.85, 0.22, 3.25, 2.03, border_col=GREEN, bw=1.5)
tb(s, 9.88, 2.32, 3.22, 0.28,
   "demo complete  |  all 28 zones warned",
   sz=8, col=GREEN, mono=True, italic=True)

# Bottom credit strip
tb(s, 0.5, 6.78, 12.0, 0.38,
   "Deploy: Vercel + Railway  |  Voice: Twilio + Amazon Polly Aditi  |  Data: NASA DONKI + NOAA SWPC + Local archive",
   sz=9, col=MUTED, mono=True, align=PP_ALIGN.CENTER)

accent_bar(s)


# ── Save ──────────────────────────────────────────────────────────────
out = r"E:\FARAWAY\docs\KAVACH_FARAWAY2026_FINAL_v4.pptx"
prs.save(out)
print(f"Saved: {len(prs.slides)} slides -> {out}")
